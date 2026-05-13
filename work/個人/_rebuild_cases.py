#!/usr/bin/env python3
"""
職務経歴書 - 事例パート再構成
- 既存 CASE STUDIES (1/3, 2/3, 3/3) を削除
- 業界別扉4枚 + 1事例1ページ × 13 を挿入
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from copy import deepcopy
from lxml import etree

SRC = '/Users/apple/.cursor/work/個人/田村太一_職務経歴書.pptx'
DST = '/Users/apple/.cursor/work/個人/田村太一_職務経歴書.pptx'

# === Color palette ===
NAVY = RGBColor(0x1B, 0x2A, 0x4A)
NAVY_LT = RGBColor(0x33, 0x4A, 0x73)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
BLUE_LT = RGBColor(0xDB, 0xEA, 0xFE)
BG_LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
SEP = RGBColor(0xE2, 0xE8, 0xF0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x10, 0xB9, 0x81)
GREEN_LT = RGBColor(0xD1, 0xFA, 0xE5)
RED = RGBColor(0xEF, 0x44, 0x44)
RED_LT = RGBColor(0xFE, 0xE2, 0xE2)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
AMBER_LT = RGBColor(0xFE, 0xF3, 0xC7)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_LT = RGBColor(0xED, 0xE9, 0xFE)
PINK = RGBColor(0xEC, 0x48, 0x99)
PINK_LT = RGBColor(0xFC, 0xE7, 0xF3)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CYAN_LT = RGBColor(0xCF, 0xFA, 0xFE)
SLATE_700 = RGBColor(0x33, 0x44, 0x55)
SLATE_500 = RGBColor(0x64, 0x74, 0x8B)
SLATE_400 = RGBColor(0x94, 0xa3, 0xb8)
SLATE_200 = RGBColor(0xE2, 0xE8, 0xF0)
FONT = 'Yu Gothic'

# Slide dimensions: 13.33in x 7.50in
SW = Inches(13.33)
SH = Inches(7.50)

# === Helpers ===
def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def no_line(shape):
    shape.line.fill.background()

def add_rect(slide, left, top, w, h, fill, line=None, line_w=None, rounded=False):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shp, left, top, w, h)
    if rounded:
        try:
            shape.adjustments[0] = 0.12
        except Exception:
            pass
    set_fill(shape, fill)
    if line is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line
        if line_w is not None:
            shape.line.width = line_w
    return shape

def add_oval(slide, left, top, w, h, fill, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    set_fill(shape, fill)
    if line is None:
        no_line(shape)
    else:
        shape.line.color.rgb = line
    return shape

def add_triangle(slide, left, top, w, h, fill, direction='right'):
    shp_map = {
        'right': MSO_SHAPE.RIGHT_TRIANGLE,
        'isos': MSO_SHAPE.ISOSCELES_TRIANGLE,
    }
    shape = slide.shapes.add_shape(shp_map.get(direction, MSO_SHAPE.ISOSCELES_TRIANGLE), left, top, w, h)
    set_fill(shape, fill)
    no_line(shape)
    return shape

def add_arrow(slide, left, top, w, h, fill, direction='right'):
    shp_map = {
        'right': MSO_SHAPE.RIGHT_ARROW,
        'left': MSO_SHAPE.LEFT_ARROW,
        'down': MSO_SHAPE.DOWN_ARROW,
        'up': MSO_SHAPE.UP_ARROW,
    }
    shape = slide.shapes.add_shape(shp_map[direction], left, top, w, h)
    set_fill(shape, fill)
    no_line(shape)
    return shape

def add_chevron(slide, left, top, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, w, h)
    set_fill(shape, fill)
    no_line(shape)
    return shape

def add_trapezoid(slide, left, top, w, h, fill, flip=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, top, w, h)
    set_fill(shape, fill)
    no_line(shape)
    if flip:
        shape.rotation = 180
    return shape

def add_line(slide, x1, y1, x2, y2, color, width=Pt(1)):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = width
    return line

def add_text(slide, left, top, w, h, text, size=11, bold=False, color=NAVY,
             align='left', anchor='top', font=FONT, line_spacing=None):
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    anchor_map = {'top': MSO_ANCHOR.TOP, 'middle': MSO_ANCHOR.MIDDLE, 'bottom': MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = anchor_map[anchor]
    align_map = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align_map[align]
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        try:
            run.font.name_far_east = font  # noqa
        except Exception:
            pass
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    # force East-Asian font via XML
    _force_east_asian_font(tb.text_frame, font)
    return tb

def _force_east_asian_font(text_frame, font_name):
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    for p in text_frame._txBody.findall('.//a:p', nsmap):
        for r in p.findall('.//a:r', nsmap):
            rPr = r.find('a:rPr', nsmap)
            if rPr is None:
                continue
            for tag in ('ea', 'cs'):
                el = rPr.find(f'a:{tag}', nsmap)
                if el is None:
                    el = etree.SubElement(rPr, f'{{{nsmap["a"]}}}{tag}')
                el.set('typeface', font_name)

def add_centered_text_in_shape(slide, shape, text, size=10, bold=False, color=WHITE, font=FONT):
    """Add a textbox centered on top of an existing shape."""
    tb = add_text(slide, shape.left, shape.top, shape.width, shape.height,
                  text, size=size, bold=bold, color=color, align='center', anchor='middle', font=font)
    return tb

# === Header / common chrome ===
def add_header(slide, page_num=None, total=13, label='CASE STUDIES'):
    add_rect(slide, 0, 0, SW, Inches(0.75), NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(8), Inches(0.4),
             label, size=14, bold=True, color=WHITE, font=FONT)
    if page_num is not None:
        add_text(slide, Inches(11.0), Inches(0.18), Inches(2.0), Inches(0.4),
                 f'{page_num:02d} / {total}', size=12, bold=True, color=WHITE, align='right', font=FONT)

def add_footer_note(slide):
    add_text(slide, Inches(0.5), Inches(7.18), Inches(12.3), Inches(0.25),
             '※ 守秘義務の関係上、社名は非公開としております。詳細は個別にご説明いたします。',
             size=8, color=SLATE_400, align='left', font=FONT)

# === Industry divider slide ===
def make_divider(slide, group_title, group_subtitle, cases_summary, accent_color, accent_lt, icon_fn=None):
    """
    group_title: '不動産'
    group_subtitle: '4 CASES'
    cases_summary: [(num, name, kpi), ...]
    icon_fn(slide, cx_in, cy_in, size_in, color) -> draws an icon
    """
    # Header
    add_rect(slide, 0, 0, SW, Inches(0.75), NAVY)
    add_text(slide, Inches(0.5), Inches(0.18), Inches(8), Inches(0.4),
             'CASE STUDIES', size=14, bold=True, color=WHITE, font=FONT)
    add_text(slide, Inches(11.0), Inches(0.18), Inches(2.0), Inches(0.4),
             f'GROUP', size=12, bold=True, color=WHITE, align='right', font=FONT)

    # Left big block with icon and group name
    add_rect(slide, 0, Inches(0.75), Inches(5.5), Inches(6.75), accent_lt)
    # accent bar
    add_rect(slide, 0, Inches(0.75), Inches(0.15), Inches(6.75), accent_color)

    if icon_fn:
        icon_fn(slide, Inches(2.75), Inches(2.4), Inches(1.8), accent_color)

    add_text(slide, Inches(0.6), Inches(4.4), Inches(4.5), Inches(0.5),
             group_subtitle, size=12, bold=True, color=accent_color, align='center', font=FONT)
    add_text(slide, Inches(0.6), Inches(4.85), Inches(4.5), Inches(1.0),
             group_title, size=44, bold=True, color=NAVY, align='center', font=FONT)
    # decorative underline
    add_rect(slide, Inches(2.5), Inches(6.05), Inches(0.7), Inches(0.05), accent_color)

    # Right: cases summary list
    rx = Inches(6.0)
    rw = Inches(6.9)
    add_text(slide, rx, Inches(1.0), rw, Inches(0.35),
             'CASES IN THIS GROUP', size=10, bold=True, color=SLATE_500, font=FONT)
    # underline
    add_rect(slide, rx, Inches(1.35), rw, Inches(0.02), accent_color)

    top = Inches(1.6)
    row_h = Inches(1.25)
    for i, (num, name, kpi) in enumerate(cases_summary):
        y = top + row_h * i
        # number badge
        add_rect(slide, rx, y, Inches(0.7), Inches(0.7), accent_color)
        add_text(slide, rx, y, Inches(0.7), Inches(0.7),
                 f'{num:02d}', size=18, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
        # case name
        add_text(slide, rx + Inches(0.85), y, rw - Inches(0.85), Inches(0.35),
                 name, size=13, bold=True, color=NAVY, font=FONT)
        # KPI highlight
        add_text(slide, rx + Inches(0.85), y + Inches(0.35), rw - Inches(0.85), Inches(0.35),
                 kpi, size=11, color=accent_color, bold=True, font=FONT)

# === Industry icons (simple geometric) ===
def icon_realestate(slide, cx, cy, size, color):
    # 家 (house)
    s = size
    # roof triangle
    roof = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, cx - s/2, cy - s/2, s, s*0.45)
    set_fill(roof, color); no_line(roof)
    # body
    body = add_rect(slide, cx - s/2 * 0.85, cy - s*0.05, s*0.85, s*0.55, color)
    # door
    door = add_rect(slide, cx - s*0.12, cy + s*0.18, s*0.24, s*0.32, RGBColor(0xFF, 0xFF, 0xFF))
    return body

def icon_btob(slide, cx, cy, size, color):
    # ビルディング (2 towers)
    s = size
    b1 = add_rect(slide, cx - s*0.45, cy - s*0.35, s*0.35, s*0.70, color)
    b2 = add_rect(slide, cx + s*0.05, cy - s*0.20, s*0.40, s*0.55, color)
    # windows on b1
    for r in range(3):
        for cc in range(2):
            w = add_rect(slide,
                         cx - s*0.40 + cc * s*0.13,
                         cy - s*0.30 + r * s*0.18,
                         s*0.08, s*0.10, WHITE)
    return b1

def icon_ec(slide, cx, cy, size, color):
    # 買い物カゴ
    s = size
    # cart body (trapezoid)
    body = add_rect(slide, cx - s*0.40, cy - s*0.10, s*0.80, s*0.40, color)
    # handle
    handle = add_rect(slide, cx - s*0.45, cy - s*0.30, s*0.20, s*0.10, color)
    # wheels
    add_oval(slide, cx - s*0.30, cy + s*0.35, s*0.15, s*0.15, color)
    add_oval(slide, cx + s*0.15, cy + s*0.35, s*0.15, s*0.15, color)
    return body

def icon_shop(slide, cx, cy, size, color):
    # 店舗 (awning + storefront)
    s = size
    # awning
    awning = add_rect(slide, cx - s*0.5, cy - s*0.35, s*1.0, s*0.18, color)
    # store body
    body = add_rect(slide, cx - s*0.45, cy - s*0.15, s*0.90, s*0.55, RGBColor(0xFF, 0xFF, 0xFF))
    body.line.color.rgb = color
    body.line.width = Pt(2.5)
    # door
    door = add_rect(slide, cx - s*0.12, cy + s*0.05, s*0.24, s*0.35, color)
    return awning

# === Case slide builder ===
def make_case_slide(slide, case):
    """
    case dict:
      num, total, industry, role, channels, challenge, approach,
      kpis [(value, label, color), ...],
      detail (extra description shown below approach),
      diagram_fn(slide, x, y, w, h)
    """
    add_header(slide, page_num=case['num'], total=case['total'])

    # Industry block
    add_text(slide, Inches(0.5), Inches(0.95), Inches(8.5), Inches(0.5),
             case['industry'], size=24, bold=True, color=NAVY, font=FONT)
    add_text(slide, Inches(0.5), Inches(1.45), Inches(8.5), Inches(0.3),
             case['role'] + '   |   ' + case['channels'], size=10, color=SLATE_500, font=FONT)
    # blue underline
    add_rect(slide, Inches(0.5), Inches(1.8), Inches(0.5), Inches(0.05), BLUE)
    add_rect(slide, Inches(1.0), Inches(1.8), Inches(11.8), Inches(0.02), SEP)

    # Left column - text content (課題, アプローチ, 成果)
    lx = Inches(0.5)
    lw = Inches(5.7)

    # Section: 課題
    add_section_header(slide, lx, Inches(2.05), '01', '課題', RED)
    add_text(slide, lx + Inches(0.15), Inches(2.55), lw - Inches(0.15), Inches(0.9),
             case['challenge'], size=11, color=SLATE_700, font=FONT, line_spacing=1.3)

    # Section: アプローチ
    add_section_header(slide, lx, Inches(3.55), '02', 'アプローチ', BLUE)
    add_text(slide, lx + Inches(0.15), Inches(4.05), lw - Inches(0.15), Inches(1.4),
             case['approach'], size=11, color=SLATE_700, font=FONT, line_spacing=1.3)

    # Section: 成果
    add_section_header(slide, lx, Inches(5.50), '03', '成果', GREEN)
    # KPI tiles
    kpi_y = Inches(6.0)
    kpis = case['kpis']
    tile_w_total = lw - Inches(0.15)
    gap = Inches(0.15)
    tile_w = (tile_w_total - gap * (len(kpis) - 1)) / len(kpis) if len(kpis) > 1 else tile_w_total
    for i, (val, lbl, col) in enumerate(kpis):
        tx = lx + Inches(0.15) + (tile_w + gap) * i
        add_rect(slide, tx, kpi_y, tile_w, Inches(1.05), GREEN_LT, rounded=True)
        add_text(slide, tx, kpi_y + Inches(0.05), tile_w, Inches(0.55),
                 val, size=22, bold=True, color=col, align='center', anchor='middle', font=FONT)
        add_text(slide, tx, kpi_y + Inches(0.60), tile_w, Inches(0.4),
                 lbl, size=9, color=SLATE_500, align='center', anchor='middle', font=FONT)

    # Right column - diagram
    rx = Inches(6.5)
    ry = Inches(2.0)
    rw = Inches(6.4)
    rh = Inches(5.0)
    # diagram background panel
    add_rect(slide, rx, ry, rw, rh, BG_LIGHT, rounded=True)
    # panel header
    add_text(slide, rx + Inches(0.3), ry + Inches(0.15), rw - Inches(0.6), Inches(0.35),
             '施策の図解', size=10, bold=True, color=SLATE_500, font=FONT)
    add_rect(slide, rx + Inches(0.3), ry + Inches(0.50), Inches(0.4), Inches(0.04), BLUE)

    # content area: between panel header (0.65 below ry) and banner (0.65 above ry+rh)
    diag_x = rx + Inches(0.3)
    diag_y = ry + Inches(0.65)
    diag_w = rw - Inches(0.6)
    diag_h = rh - Inches(0.65) - Inches(0.65)  # leave room for banner
    case['diagram_fn'](slide, diag_x, diag_y, diag_w, diag_h)

    # Common banner across all diagrams
    banner_top = ry + rh - Inches(0.65)
    add_rect(slide, diag_x, banner_top, diag_w, Inches(0.55), NAVY, rounded=True)
    add_text(slide, diag_x, banner_top, diag_w, Inches(0.55),
             case.get('banner', ''), size=10, bold=True, color=WHITE,
             align='center', anchor='middle', font=FONT)

    add_footer_note(slide)

def add_section_header(slide, x, y, num, label, color):
    # number circle
    add_oval(slide, x, y, Inches(0.30), Inches(0.30), color)
    add_text(slide, x, y, Inches(0.30), Inches(0.30),
             num, size=9, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
    add_text(slide, x + Inches(0.40), y - Inches(0.02), Inches(4.0), Inches(0.35),
             label, size=14, bold=True, color=NAVY, font=FONT)
    # accent line right
    add_rect(slide, x + Inches(0.40), y + Inches(0.30), Inches(5.0), Inches(0.02), SEP)

# === Diagram functions ===
# Each takes (slide, x, y, w, h)

def diag_box(slide, x, y, w, h, text, fill, color=WHITE, size=10, bold=True, rounded=True):
    s = add_rect(slide, x, y, w, h, fill, rounded=rounded)
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=color,
             align='center', anchor='middle', font=FONT)
    return s

def diag_outlined_box(slide, x, y, w, h, text, line_color, fill=WHITE, color=NAVY, size=10, bold=False):
    s = add_rect(slide, x, y, w, h, fill, rounded=True, line=line_color, line_w=Pt(1.5))
    add_text(slide, x, y, w, h, text, size=size, bold=bold, color=color,
             align='center', anchor='middle', font=FONT)
    return s

def diag_arrow_right(slide, x, y, w, h, color=SLATE_400):
    add_arrow(slide, x, y, w, h, color, direction='right')

def diag_arrow_down(slide, x, y, w, h, color=SLATE_400):
    add_arrow(slide, x, y, w, h, color, direction='down')


# ----- CASE 01: 不動産（売買領域）大規模運用 -----
def diag_case01(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '準顕在層向け動画 × 顧客データ機械学習最適化',
             size=11, bold=True, color=NAVY, align='center', font=FONT)
    fx = x + Inches(0.2)
    fy = y + Inches(0.45)
    # funnel stages (3 levels)
    add_trapezoid(slide, fx, fy, Inches(2.8), Inches(0.6), BLUE_LT, flip=True)
    add_text(slide, fx, fy, Inches(2.8), Inches(0.6), '潜在層', size=11, bold=True, color=NAVY, align='center', anchor='middle', font=FONT)
    add_trapezoid(slide, fx + Inches(0.3), fy + Inches(0.75), Inches(2.2), Inches(0.6), AMBER_LT, flip=True)
    add_text(slide, fx + Inches(0.3), fy + Inches(0.75), Inches(2.2), Inches(0.6), '準顕在層', size=11, bold=True, color=AMBER, align='center', anchor='middle', font=FONT)
    add_trapezoid(slide, fx + Inches(0.6), fy + Inches(1.50), Inches(1.6), Inches(0.6), GREEN_LT, flip=True)
    add_text(slide, fx + Inches(0.6), fy + Inches(1.50), Inches(1.6), Inches(0.6), 'CV（問合せ）', size=11, bold=True, color=GREEN, align='center', anchor='middle', font=FONT)

    # right side - 3 boxes for new strategies
    rx = x + Inches(3.4)
    rw2 = w - Inches(3.4) - Inches(0.1)
    add_text(slide, rx, fy - Inches(0.15), rw2, Inches(0.3), '⚡ 新規施策', size=10, bold=True, color=BLUE, font=FONT)
    diag_box(slide, rx, fy + Inches(0.15), rw2, Inches(0.55),
             '🎬  準顕在向け動画施策', BLUE, size=10, bold=True)
    diag_box(slide, rx, fy + Inches(0.78), rw2, Inches(0.55),
             '🤖  顧客データ × ML 最適化', PURPLE, size=10, bold=True)
    diag_box(slide, rx, fy + Inches(1.41), rw2, Inches(0.55),
             '🎯  媒体ミックス再設計', NAVY, size=10, bold=True)
    # arrow from new strategies into the funnel's middle stage
    add_arrow(slide, x + Inches(3.05), fy + Inches(0.9), Inches(0.35), Inches(0.25), BLUE, direction='left')

    # Scale info row at the bottom of content area
    sy = y + Inches(2.7)
    add_rect(slide, x + Inches(0.1), sy, w - Inches(0.2), Inches(0.85), BG_LIGHT, rounded=True, line=SLATE_200, line_w=Pt(0.75))
    add_text(slide, x + Inches(0.2), sy + Inches(0.10), w - Inches(0.4), Inches(0.30),
             'スケール感', size=9, bold=True, color=SLATE_500, align='left', font=FONT)
    add_text(slide, x + Inches(0.2), sy + Inches(0.40), w - Inches(0.4), Inches(0.35),
             '月間予算  ¥2億規模  ／  並行運用  5媒体  ／  運用期間  数年スパン',
             size=11, bold=True, color=NAVY, align='center', anchor='middle', font=FONT)

# ----- CASE 02: 不動産（法人向け売買）為替×リード相関 -----
def diag_case02(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '為替指標連動 × 運用額の動的調整ロジック',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    # 左: 発見した相関
    lx = x + Inches(0.1)
    lw = Inches(2.6)
    add_text(slide, lx, y + Inches(0.4), lw, Inches(0.3), '💡 独自発見の相関', size=10, bold=True, color=AMBER, font=FONT)
    add_rect(slide, lx, y + Inches(0.7), lw, Inches(2.5), WHITE, rounded=True, line=SLATE_200, line_w=Pt(1))
    add_text(slide, lx + Inches(0.15), y + Inches(0.8), lw - Inches(0.3), Inches(0.25),
             '為替（¥/$）', size=9, bold=True, color=SLATE_700, font=FONT)
    bx = lx + Inches(0.25)
    by = y + Inches(1.65)
    heights = [0.25, 0.45, 0.35, 0.65, 0.85]
    for i, hv in enumerate(heights):
        add_rect(slide, bx + Emu(int(0.40 * i * 914400)), by - Inches(hv), Inches(0.28), Inches(hv), BLUE)
    add_text(slide, lx + Inches(0.15), y + Inches(1.80), lw - Inches(0.3), Inches(0.25),
             'リード数', size=9, bold=True, color=GREEN, font=FONT, align='left')
    # mock リード count points overlaying
    for i, hv in enumerate(heights):
        add_oval(slide, bx + Emu(int(0.40 * i * 914400)) + Inches(0.08), by - Inches(hv) - Inches(0.10), Inches(0.13), Inches(0.13), GREEN)
    add_text(slide, lx + Inches(0.15), y + Inches(2.75), lw - Inches(0.3), Inches(0.3),
             '為替↑  ⇒  リード↑', size=11, bold=True, color=RED, align='center', font=FONT)

    # 矢印
    add_arrow(slide, x + Inches(2.78), y + Inches(1.65), Inches(0.32), Inches(0.4), BLUE)

    # 右: 動的調整ロジック
    rx = x + Inches(3.20)
    rw = w - Inches(3.30)
    add_text(slide, rx, y + Inches(0.4), rw, Inches(0.3), '⚙ 動的調整ロジック', size=10, bold=True, color=BLUE, font=FONT)
    add_rect(slide, rx, y + Inches(0.7), rw, Inches(2.5), WHITE, rounded=True, line=SLATE_200, line_w=Pt(1))

    rows = [
        ('為替 高水準', '運用額 拡大', GREEN),
        ('為替 中水準', '運用額 維持', AMBER),
        ('為替 低水準', '運用額 抑制', RED),
    ]
    cond_w = Inches(1.10)
    arrow_w = Inches(0.28)
    act_x_offset = Inches(0.12) + cond_w + Inches(0.04) + arrow_w + Inches(0.04)
    act_w = rw - act_x_offset - Inches(0.12)
    rty = y + Inches(0.85)
    for i, (cond, act, col) in enumerate(rows):
        ty = rty + Inches(0.72 * i)
        diag_box(slide, rx + Inches(0.12), ty, cond_w, Inches(0.50), cond, col, size=9, bold=True)
        add_arrow(slide, rx + Inches(0.12) + cond_w + Inches(0.04), ty + Inches(0.12), arrow_w, Inches(0.26), SLATE_400)
        diag_outlined_box(slide, rx + act_x_offset, ty, act_w, Inches(0.50), act, col, color=col, size=10, bold=True)

# ----- CASE 03: 不動産（賃貸管理）マッチタイプ最適化 -----
def diag_case03(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), 'マッチタイプ × 最適化ポイントの戦略的見直し',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x, y + Inches(0.35), w, Inches(0.3), '🎯 マッチタイプ設計', size=10, bold=True, color=BLUE, font=FONT)
    mt = [
        ('完全一致', '主要KW\n（高効率）', GREEN),
        ('フレーズ', '購入意向KW\n（拡張）', AMBER),
        ('部分一致', '探索 + Smart Bid\n（学習）', PURPLE),
    ]
    tile_w = (w - Inches(0.4)) / 3
    for i, (t, d, c) in enumerate(mt):
        tx = x + Inches(0.1) + (tile_w + Inches(0.1)) * i
        diag_box(slide, tx, y + Inches(0.7), tile_w, Inches(0.45), t, c, size=11, bold=True)
        add_rect(slide, tx, y + Inches(1.15), tile_w, Inches(0.70), WHITE, rounded=True, line=SLATE_200, line_w=Pt(1))
        add_text(slide, tx, y + Inches(1.15), tile_w, Inches(0.70), d,
                 size=10, color=SLATE_700, align='center', anchor='middle', font=FONT)

    for i in range(3):
        tx = x + Inches(0.1) + (tile_w + Inches(0.1)) * i + tile_w/2 - Inches(0.12)
        add_arrow(slide, tx, y + Inches(1.95), Inches(0.24), Inches(0.20), SLATE_400, direction='down')

    add_rect(slide, x + Inches(0.1), y + Inches(2.25), w - Inches(0.2), Inches(1.30), BLUE_LT, rounded=True)
    add_text(slide, x + Inches(0.1), y + Inches(2.30), w - Inches(0.2), Inches(0.3),
             '⚙ 最適化ポイントの戦略再設計', size=10, bold=True, color=NAVY, align='center', font=FONT)
    points = ['ターゲットCPA', 'コンバージョン値', '入札戦略の段階移行']
    pw = (w - Inches(0.6)) / 3
    for i, p in enumerate(points):
        px = x + Inches(0.2) + (pw + Inches(0.1)) * i
        diag_outlined_box(slide, px, y + Inches(2.70), pw, Inches(0.75), p, BLUE, color=NAVY, size=10, bold=True)

# ----- CASE 04: 不動産（地域特化 戸建て販売）全チャネル統合 -----
def diag_case04(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), 'マーケティング全チャネルの統合ディレクション',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    cx = x + w/2
    cy = y + Inches(1.95)
    hub_size = Inches(1.25)
    hub_x = cx - hub_size/2
    hub_y = cy - hub_size/2
    add_oval(slide, hub_x, hub_y, hub_size, hub_size, NAVY)

    # 6 satellites around the hub (compressed orbit to fit h≈3.7)
    sats = [
        ('広告', BLUE, -2.20, -1.20),
        ('SNS', PURPLE, 0, -1.40),
        ('LINE', GREEN, 2.20, -1.20),
        ('IG', PINK, 2.20, 1.05),
        ('インフルエンサー', AMBER, 0, 1.30),
        ('WP制作', CYAN, -2.20, 1.05),
    ]
    sat_w = Inches(1.4)
    sat_h = Inches(0.50)
    # draw connections first
    for name, col, dx, dy in sats:
        sx = cx + Inches(dx) - sat_w/2
        sy = cy + Inches(dy) - sat_h/2
        add_line(slide, cx, cy, sx + sat_w/2, sy + sat_h/2, SLATE_400, width=Pt(0.75))

    # redraw hub on top of lines
    add_oval(slide, hub_x, hub_y, hub_size, hub_size, NAVY)
    add_text(slide, hub_x, hub_y, hub_size, hub_size,
             '統合\nディレクション', size=11, bold=True, color=WHITE,
             align='center', anchor='middle', font=FONT)

    # draw satellites on top
    for name, col, dx, dy in sats:
        sx = cx + Inches(dx) - sat_w/2
        sy = cy + Inches(dy) - sat_h/2
        diag_box(slide, sx, sy, sat_w, sat_h, name, col, size=10, bold=True)

# ----- CASE 05: BtoB（知財系）検索市場特性分析 -----
def diag_case05(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '業界特有の検索市場構造を深く分析・運用設計に反映',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '🔎 検索クエリ階層マップ', size=10, bold=True, color=BLUE, font=FONT)

    levels = [
        ('Tier 1 :  業界特化用語（指名級）', '高CVR・低ボリューム', GREEN, 0.0),
        ('Tier 2 :  ニーズ顕在クエリ',       '中CVR・中ボリューム', AMBER, 0.4),
        ('Tier 3 :  関連業務クエリ',          'CVR要設計・大ボリューム', BLUE, 0.8),
    ]
    ly = y + Inches(0.65)
    for i, (label, desc, col, indent) in enumerate(levels):
        ty = ly + Inches(0.52 * i)
        add_rect(slide, x + Inches(0.2) + Inches(indent), ty, Inches(3.4) - Inches(indent), Inches(0.42), col, rounded=True)
        add_text(slide, x + Inches(0.2) + Inches(indent), ty, Inches(3.4) - Inches(indent), Inches(0.42),
                 label, size=10, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
        add_text(slide, x + Inches(3.7), ty + Inches(0.06), w - Inches(3.8), Inches(0.32),
                 desc, size=9, color=SLATE_700, font=FONT)

    add_rect(slide, x + Inches(0.1), y + Inches(2.35), w - Inches(0.2), Inches(1.20), AMBER_LT, rounded=True)
    add_text(slide, x + Inches(0.1), y + Inches(2.42), w - Inches(0.2), Inches(0.3),
             '⚙ 市場特性を最大限活かした運用設計', size=10, bold=True, color=AMBER, align='center', font=FONT)
    items = ['Tier別キャンペーン分離', 'CV値の重み付け', '社内ナレッジ蓄積']
    iw = (w - Inches(0.6)) / 3
    for i, it in enumerate(items):
        ix = x + Inches(0.2) + (iw + Inches(0.1)) * i
        diag_outlined_box(slide, ix, y + Inches(2.78), iw, Inches(0.62), it, AMBER, color=NAVY, size=10, bold=True)

# ----- CASE 06: 人材（特化型転職）LP×戦略リプレイス -----
def diag_case06(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '戦略設計から実行まで一気通貫  |  LP改善で CVR 125%',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '📉 リプレイス3ヶ月で 登録CPA 約50%削減', size=10, bold=True, color=GREEN, font=FONT)

    # Before/After bars
    bar_x = x + Inches(0.95)
    bar_full_w = w - Inches(1.05) - Inches(0.10)  # max width
    bar_h = Inches(0.45)
    # Before
    bar_y = y + Inches(0.70)
    add_text(slide, x + Inches(0.1), bar_y + Inches(0.10), Inches(0.8), Inches(0.3), 'Before', size=10, bold=True, color=SLATE_500, font=FONT)
    add_rect(slide, bar_x, bar_y, bar_full_w, bar_h, RED)
    add_text(slide, bar_x, bar_y, bar_full_w, bar_h, '旧体制 CPA（100%）',
             size=10, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
    # After
    bar_y2 = y + Inches(1.30)
    add_text(slide, x + Inches(0.1), bar_y2 + Inches(0.10), Inches(0.8), Inches(0.3), 'After', size=10, bold=True, color=GREEN, font=FONT)
    after_w = bar_full_w * 0.5
    add_rect(slide, bar_x, bar_y2, after_w, bar_h, GREEN)
    add_text(slide, bar_x, bar_y2, after_w, bar_h, '50%削減',
             size=10, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
    add_arrow(slide, bar_x + after_w + Inches(0.05), bar_y2 + Inches(0.10), Inches(0.45), Inches(0.25), AMBER, direction='right')
    add_text(slide, bar_x + after_w + Inches(0.55), bar_y2, bar_full_w - after_w - Inches(0.55), bar_h,
             '3ヶ月で達成', size=10, bold=True, color=AMBER, anchor='middle', font=FONT)

    # フロー
    add_text(slide, x + Inches(0.1), y + Inches(2.05), w - Inches(0.2), Inches(0.3),
             '🛠 戦略設計 → LP改善 → 一気通貫実行', size=10, bold=True, color=BLUE, font=FONT)
    flow = ['戦略\n再設計', '自社アセット\n活用LP', 'LPO\nABテスト', 'CVR\n125%']
    colors = [NAVY, BLUE, PURPLE, GREEN]
    n = len(flow)
    arrow_w = Inches(0.25)
    gap = Inches(0.10)
    box_w = (w - Inches(0.2) - arrow_w * (n - 1) - gap * 2 * (n - 1)) / n
    box_h = Inches(0.75)
    flow_y = y + Inches(2.45)
    cur_x = x + Inches(0.1)
    for i, (txt, col) in enumerate(zip(flow, colors)):
        diag_box(slide, cur_x, flow_y, box_w, box_h, txt, col, size=10, bold=True)
        cur_x = cur_x + box_w
        if i < n - 1:
            add_arrow(slide, cur_x + gap, flow_y + Inches(0.25), arrow_w, Inches(0.25), SLATE_400)
            cur_x = cur_x + arrow_w + gap * 2

# ----- CASE 07: BtoB（SNS拡張ツール）データ横断 -----
def diag_case07(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '顧客・商談・SFデータの横断分析  ⇒  運用設計を根本再構築',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '🔗 3つのデータソースを統合分析', size=10, bold=True, color=BLUE, font=FONT)

    sources = [
        ('🎤 顧客ヒアリング', '定性ニーズ・課題', PURPLE),
        ('💬 商談データ', '受注/失注理由', BLUE),
        ('📊 SF（Salesforce）', 'パイプライン全量', GREEN),
    ]
    sw = (w - Inches(0.5)) / 3
    sy = y + Inches(0.70)
    for i, (label, desc, col) in enumerate(sources):
        sx = x + Inches(0.15) + (sw + Inches(0.1)) * i
        add_rect(slide, sx, sy, sw, Inches(0.85), col, rounded=True)
        add_text(slide, sx, sy + Inches(0.08), sw, Inches(0.32), label,
                 size=11, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
        add_text(slide, sx, sy + Inches(0.45), sw, Inches(0.35), desc,
                 size=9, color=WHITE, align='center', anchor='middle', font=FONT)

    # converging arrows
    for i in range(3):
        sx = x + Inches(0.15) + (sw + Inches(0.1)) * i + sw/2 - Inches(0.12)
        add_arrow(slide, sx, y + Inches(1.65), Inches(0.24), Inches(0.24), SLATE_400, direction='down')

    add_rect(slide, x + Inches(0.3), y + Inches(2.00), w - Inches(0.6), Inches(0.50), NAVY, rounded=True)
    add_text(slide, x + Inches(0.3), y + Inches(2.00), w - Inches(0.6), Inches(0.50),
             '横断分析  ⇒  運用設計の根本再構築',
             size=11, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)

    add_arrow(slide, x + w/2 - Inches(0.12), y + Inches(2.58), Inches(0.24), Inches(0.22), SLATE_400, direction='down')

    outcomes = [('運用設計\n刷新', BLUE), ('サイト改善\nLPO', PURPLE), ('商談 1.5倍\nCVR 135%', GREEN)]
    oy = y + Inches(2.90)
    ow = (w - Inches(0.5)) / 3
    for i, (t, c) in enumerate(outcomes):
        ox = x + Inches(0.15) + (ow + Inches(0.1)) * i
        diag_box(slide, ox, oy, ow, Inches(0.65), t, c, size=10, bold=True)

# ----- CASE 08: BtoB（採用サービス）営業代行統合 -----
def diag_case08(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '広告領域を超えた営業体制までの統合施策',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    panel_h = Inches(3.10)
    panel_y = y + Inches(0.40)
    lx = x + Inches(0.1)
    lw = (w - Inches(0.5)) / 2
    add_rect(slide, lx, panel_y, lw, panel_h, BLUE_LT, rounded=True)
    add_text(slide, lx, panel_y + Inches(0.05), lw, Inches(0.3), '【 広告施策 】', size=11, bold=True, color=BLUE, align='center', font=FONT)
    ad_items = ['リスティング', 'ディスプレイ', 'SNS広告', 'DM施策', 'AI活用クリエイティブ']
    for i, item in enumerate(ad_items):
        diag_outlined_box(slide, lx + Inches(0.3), panel_y + Inches(0.45) + Inches(0.46 * i),
                          lw - Inches(0.6), Inches(0.36), item, BLUE, color=NAVY, size=10, bold=False)

    rx = lx + lw + Inches(0.3)
    add_rect(slide, rx, panel_y, lw, panel_h, AMBER_LT, rounded=True)
    add_text(slide, rx, panel_y + Inches(0.05), lw, Inches(0.3), '【 営業代行・支援 】', size=11, bold=True, color=AMBER, align='center', font=FONT)
    add_text(slide, rx, panel_y + Inches(0.35), lw, Inches(0.3), '広告領域を超えた施策', size=9, color=SLATE_500, align='center', font=FONT)
    diag_box(slide, rx + Inches(0.3), panel_y + Inches(0.75), lw - Inches(0.6), Inches(0.50),
             '営業ボトルネック特定', AMBER, size=10, bold=True)
    add_arrow(slide, rx + lw/2 - Inches(0.13), panel_y + Inches(1.30), Inches(0.26), Inches(0.22), SLATE_400, direction='down')
    diag_box(slide, rx + Inches(0.3), panel_y + Inches(1.58), lw - Inches(0.6), Inches(0.50),
             '営業代行 企画・実行', AMBER, size=10, bold=True)
    add_arrow(slide, rx + lw/2 - Inches(0.13), panel_y + Inches(2.13), Inches(0.26), Inches(0.22), SLATE_400, direction='down')
    diag_box(slide, rx + Inches(0.3), panel_y + Inches(2.40), lw - Inches(0.6), Inches(0.55),
             '商談 +50件/月', GREEN, size=11, bold=True)

# ----- CASE 09: EC（住宅設備）フルファネル -----
def diag_case09(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), 'フルファネル施策 × 媒体マッピング',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    # 4段ファネル（全幅・降順）
    stages = [
        ('認知', '広く想起獲得', '認知広告 / Instagram', PURPLE,  5.40),
        ('興味', '商品理解促進', 'ディスプレイ / SNS',   BLUE,    4.60),
        ('比較', '購入候補化',   'リスティング',          AMBER,   3.80),
        ('購入', 'CV刈り取り',   'ダイナミック広告',      GREEN,   3.00),
    ]
    sh = Inches(0.48)
    gap_y = Inches(0.10)
    sy = y + Inches(0.45)
    for i, (label, desc, ch, col, ws_in) in enumerate(stages):
        stage_w = Inches(ws_in)
        sx = x + (w - stage_w) / 2
        add_rect(slide, sx, sy, stage_w, sh, col, rounded=True)
        add_text(slide, sx + Inches(0.10), sy, Inches(0.7), sh, label,
                 size=11, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
        add_text(slide, sx + Inches(0.80), sy, stage_w - Inches(1.6), sh, desc,
                 size=10, color=WHITE, align='left', anchor='middle', font=FONT)
        # channel inside stage (right side)
        add_text(slide, sx + stage_w - Inches(2.20), sy, Inches(2.10), sh, ch,
                 size=9, color=WHITE, bold=True, align='right', anchor='middle', font=FONT)
        sy = sy + sh + gap_y

    # 下: 売上260% + 運用費4倍 を並べる
    box_y = y + Inches(3.10)
    box_h = Inches(0.55)
    half_w = (w - Inches(0.3)) / 2
    add_rect(slide, x + Inches(0.1), box_y, half_w, box_h, GREEN_LT, rounded=True, line=GREEN, line_w=Pt(1.5))
    add_text(slide, x + Inches(0.1), box_y, half_w, box_h,
             '🚀  売上  260%  に伸長',
             size=12, bold=True, color=GREEN, align='center', anchor='middle', font=FONT)
    add_rect(slide, x + Inches(0.2) + half_w, box_y, half_w, box_h, AMBER_LT, rounded=True, line=AMBER, line_w=Pt(1.5))
    add_text(slide, x + Inches(0.2) + half_w, box_y, half_w, box_h,
             '💰  運用費  約4倍  に拡大獲得',
             size=12, bold=True, color=AMBER, align='center', anchor='middle', font=FONT)

# ----- CASE 10: 店舗系（車整備）地域特性データ -----
def diag_case10(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '地域特性データを独自調査  ⇒  エリア別配信戦略を精緻設計',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    # 上: 4つのデータソース（横並び）
    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '📍 独自調査した 4種類の地域特性データ', size=10, bold=True, color=BLUE, font=FONT)
    data_items = [
        ('🚗\n車保有人口', BLUE),
        ('🛣\n高速道路有無', PURPLE),
        ('⛰\n店舗周囲の地形', GREEN),
        ('🏘\n人口密度', AMBER),
    ]
    iw = (w - Inches(0.4)) / 4
    for i, (it, col) in enumerate(data_items):
        ix_ = x + Inches(0.1) + (iw + Inches(0.067)) * i
        diag_box(slide, ix_, y + Inches(0.70), iw, Inches(0.75), it, col, size=10, bold=True)

    # 中央矢印
    add_arrow(slide, x + w/2 - Inches(0.13), y + Inches(1.55), Inches(0.26), Inches(0.22), SLATE_400, direction='down')

    # 下: 左ヒートマップ、右 配信戦略 3段
    map_x = x + Inches(0.1)
    map_w = Inches(2.95)
    map_y = y + Inches(1.85)
    map_h = Inches(1.85)
    add_rect(slide, map_x, map_y, map_w, map_h, WHITE, rounded=True, line=SLATE_200, line_w=Pt(1))
    add_text(slide, map_x, map_y + Inches(0.05), map_w, Inches(0.30), '配信ヒートマップ',
             size=10, bold=True, color=SLATE_500, align='center', font=FONT)
    heat_data = [
        (0.25, 0.50, RED, 0.40),
        (0.55, 0.45, AMBER, 0.32),
        (0.75, 0.55, GREEN, 0.28),
        (0.35, 0.70, RED, 0.35),
        (0.65, 0.78, AMBER, 0.28),
        (0.20, 0.82, BLUE, 0.18),
    ]
    for fx, fy, col, sz in heat_data:
        cx = map_x + Inches(map_w.inches * fx) - Inches(sz/2)
        cy = map_y + Inches(map_h.inches * fy) - Inches(sz/2)
        add_oval(slide, cx, cy, Inches(sz), Inches(sz), col)

    # arrow to strategies
    add_arrow(slide, map_x + map_w + Inches(0.03), y + Inches(2.65), Inches(0.30), Inches(0.22), SLATE_400)

    # 右: 配信戦略
    sx_ = map_x + map_w + Inches(0.40)
    sw_ = x + w - sx_ - Inches(0.1)
    add_text(slide, sx_, map_y + Inches(0.0), sw_, Inches(0.30),
             '🎯 エリア別 配信戦略', size=10, bold=True, color=BLUE, font=FONT)
    strats = [
        ('高密度エリア', '入札強化 / 高頻度', RED),
        ('中密度エリア', '効率重視', AMBER),
        ('低密度エリア', '抑制 / 除外', SLATE_500),
    ]
    for i, (name, act, col) in enumerate(strats):
        ty = map_y + Inches(0.35) + Inches(0.52 * i)
        add_rect(slide, sx_, ty, sw_, Inches(0.45), WHITE, rounded=True, line=col, line_w=Pt(1.5))
        add_text(slide, sx_ + Inches(0.12), ty + Inches(0.04), sw_ - Inches(0.2), Inches(0.20),
                 name, size=10, bold=True, color=col, font=FONT)
        add_text(slide, sx_ + Inches(0.12), ty + Inches(0.22), sw_ - Inches(0.2), Inches(0.20),
                 act, size=9, color=SLATE_700, font=FONT)

# ----- CASE 11: 店舗系（歯科矯正）AI動画 -----
def diag_case11(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), 'AI技術で動画モデル不足を突破  |  制作ボトルネック解消',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '⚠ Before : クリエイティブ制作のボトルネック', size=10, bold=True, color=RED, font=FONT)
    add_rect(slide, x + Inches(0.1), y + Inches(0.65), w - Inches(0.2), Inches(0.55), RED_LT, rounded=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.65), w - Inches(0.4), Inches(0.55),
             '出演モデル不足  →  動画クリエイティブを量産できない  →  改善が頭打ち',
             size=11, bold=True, color=RED, align='center', anchor='middle', font=FONT)

    add_arrow(slide, x + w/2 - Inches(0.17), y + Inches(1.25), Inches(0.34), Inches(0.30), AMBER, direction='down')

    add_text(slide, x + Inches(0.1), y + Inches(1.65), w - Inches(0.2), Inches(0.3),
             '🤖 After : AI技術で突破', size=10, bold=True, color=PURPLE, font=FONT)
    flow = [
        ('AI モデル\n生成', PURPLE),
        ('動画\nクリエイティブ', BLUE),
        ('量産 × 高速\nABテスト', AMBER),
        ('予約CPA 1/3\n来院CPA -50%', GREEN),
    ]
    fy = y + Inches(2.05)
    n = 4
    arrow_w = Inches(0.25)
    gap = Inches(0.10)
    fw = (w - Inches(0.2) - arrow_w * (n - 1) - gap * 2 * (n - 1)) / n
    fh = Inches(1.10)
    cur_x = x + Inches(0.1)
    for i, (txt, col) in enumerate(flow):
        diag_box(slide, cur_x, fy, fw, fh, txt, col, size=10, bold=True)
        cur_x = cur_x + fw
        if i < n - 1:
            add_arrow(slide, cur_x + gap, fy + Inches(0.40), arrow_w, Inches(0.30), SLATE_400)
            cur_x = cur_x + arrow_w + gap * 2

# ----- CASE 12: 店舗系（着物）定量×定性 LP改善 -----
def diag_case12(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '定量 × 定性データ で LP改善の仮説検証サイクルを設計',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '📊 データドリブンな改善サイクル', size=10, bold=True, color=BLUE, font=FONT)

    qx = x + Inches(0.1)
    qw = (w - Inches(0.4)) / 2
    add_rect(slide, qx, y + Inches(0.70), qw, Inches(0.85), BLUE_LT, rounded=True)
    add_text(slide, qx, y + Inches(0.75), qw, Inches(0.30), '【 定量データ 】', size=11, bold=True, color=BLUE, align='center', font=FONT)
    add_text(slide, qx, y + Inches(1.05), qw, Inches(0.50),
             'GA / 広告管理画面 / CV数 / CVR / 直帰率', size=10, color=SLATE_700, align='center', anchor='middle', font=FONT)

    qrx = qx + qw + Inches(0.2)
    add_rect(slide, qrx, y + Inches(0.70), qw, Inches(0.85), PURPLE_LT, rounded=True)
    add_text(slide, qrx, y + Inches(0.75), qw, Inches(0.30), '【 定性データ 】', size=11, bold=True, color=PURPLE, align='center', font=FONT)
    add_text(slide, qrx, y + Inches(1.05), qw, Inches(0.50),
             'ヒートマップ / セッションリプレイ / ユーザー行動', size=10, color=SLATE_700, align='center', anchor='middle', font=FONT)

    for fxx in [qx + qw/2 - Inches(0.12), qrx + qw/2 - Inches(0.12)]:
        add_arrow(slide, fxx, y + Inches(1.62), Inches(0.24), Inches(0.22), SLATE_400, direction='down')

    cy = y + Inches(1.95)
    cycle_w = w - Inches(0.4)
    add_rect(slide, x + Inches(0.2), cy, cycle_w, Inches(1.65), AMBER_LT, rounded=True)
    add_text(slide, x + Inches(0.2), cy + Inches(0.05), cycle_w, Inches(0.30),
             '🔄 仮説検証サイクル', size=10, bold=True, color=AMBER, align='center', font=FONT)
    cycle_items = [('仮説', AMBER), ('AB\nテスト', BLUE), ('検証', PURPLE), ('改善', GREEN)]
    n = 4
    arrow_w = Inches(0.22)
    gap = Inches(0.08)
    inner_pad = Inches(0.30)
    ciw = (cycle_w - inner_pad * 2 - arrow_w * (n - 1) - gap * 2 * (n - 1)) / n
    cur = x + Inches(0.2) + inner_pad
    for i, (t, c) in enumerate(cycle_items):
        diag_box(slide, cur, cy + Inches(0.40), ciw, Inches(0.80), t, c, size=11, bold=True)
        cur = cur + ciw
        if i < n - 1:
            add_arrow(slide, cur + gap, cy + Inches(0.65), arrow_w, Inches(0.30), SLATE_400)
            cur = cur + arrow_w + gap * 2
    # cyclic arrow back
    add_arrow(slide, x + Inches(0.4), cy + Inches(1.30), cycle_w - Inches(0.4) - inner_pad, Inches(0.20), AMBER, direction='left')

# ----- CASE 13: 買取系（商業車）多軸最適化 -----
def diag_case13(slide, x, y, w, h):
    add_text(slide, x, y, w, Inches(0.3), '市場・商材分析 × 多軸最適化マトリクス',
             size=11, bold=True, color=NAVY, align='center', font=FONT)

    add_text(slide, x + Inches(0.1), y + Inches(0.35), w - Inches(0.2), Inches(0.3),
             '⚙ 多角的に最適化した 3軸', size=10, bold=True, color=BLUE, font=FONT)

    axes = [
        ('🎨 クリエイティブ', '商材特性に合わせた\n訴求軸 × 画像/コピー', PURPLE),
        ('🗺 地域',           '商業車需要の高い\nエリア別 配信制御',   BLUE),
        ('📱 デバイス',       'デバイス別CVR・\n入札比率の最適化',   GREEN),
    ]
    aw = (w - Inches(0.4)) / 3
    ay = y + Inches(0.70)
    for i, (name, desc, col) in enumerate(axes):
        ax_ = x + Inches(0.1) + (aw + Inches(0.1)) * i
        add_rect(slide, ax_, ay, aw, Inches(1.30), col, rounded=True)
        add_text(slide, ax_, ay + Inches(0.10), aw, Inches(0.40),
                 name, size=12, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)
        add_text(slide, ax_, ay + Inches(0.55), aw, Inches(0.70),
                 desc, size=10, color=WHITE, align='center', anchor='middle', font=FONT)

    for i in range(3):
        ax_ = x + Inches(0.1) + (aw + Inches(0.1)) * i + aw/2 - Inches(0.13)
        add_arrow(slide, ax_, ay + Inches(1.35), Inches(0.26), Inches(0.22), SLATE_400, direction='down')

    cy_ = ay + Inches(1.65)
    add_rect(slide, x + Inches(0.2), cy_, w - Inches(0.4), Inches(0.50), NAVY, rounded=True)
    add_text(slide, x + Inches(0.2), cy_, w - Inches(0.4), Inches(0.50),
             '問合せ  →  買取  の転換率向上を主目的に統合最適化',
             size=11, bold=True, color=WHITE, align='center', anchor='middle', font=FONT)


# === Case data ===
CASES = [
    # 不動産グループ
    dict(
        num=1, total=13, industry='不動産（売買領域）',
        role='運用担当（2名体制）  |  月間予算 2億円規模',
        channels='リスティング / DSP / SNS / ダイナミック / 認知広告',
        challenge='大規模アカウントで既存媒体の改善余地が限定的。\n更なる伸長には新たな打ち手の創出が必要だった。',
        approach='準顕在層向けの動画施策を新規企画。あわせて顧客データを活用した機械学習最適化を導入し、媒体ミックスを再設計。広告主に提案・自ら実行までを完遂。',
        kpis=[('+130%', '問合せ数', GREEN), ('¥2億', '月間運用予算', NAVY)],
        banner='月間運用予算 2億円規模  |  メイン媒体5種（リスティング / DSP / SNS / DA / 認知）',
        diagram_fn=diag_case01,
    ),
    dict(
        num=2, total=13, industry='不動産（法人向け売買領域）',
        role='運用担当  |  アカウント新規立ち上げ',
        channels='リスティング / ディスプレイ',
        challenge='法人向け不動産は購入意向の変動要因が複雑。一般的な運用設計だけでは安定した投資効率を出しにくかった。',
        approach='為替変動とリード数の相関関係を独自に発見。為替指標に基づき運用額を動的に調整するロジックを構築し、投資効率を最大化。',
        kpis=[('+140%', '問合せ数', GREEN), ('独自', '相関分析ロジック', AMBER)],
        banner='法人不動産 × 為替感応度の高い顧客特性を独自に発見し、投資効率を最大化',
        diagram_fn=diag_case02,
    ),
    dict(
        num=3, total=13, industry='不動産（賃貸管理領域）',
        role='運用担当  |  アカウント新規立ち上げ',
        channels='リスティング / SNS広告',
        challenge='ゼロからの新規アカウント構築。短期間で成果を出す必要があり、設計の精度が成否を分ける状況。',
        approach='マッチタイプの設計と最適化ポイントの戦略的見直しを実施。短期間で成果を創出できる構造を最初から組み込み。',
        kpis=[('+120%', '問合せ数', GREEN), ('短期', '立ち上げ期間', BLUE)],
        banner='ゼロからのアカウント設計  |  短期間で成果創出  |  目標・昨対比 120%以上',
        diagram_fn=diag_case03,
    ),
    dict(
        num=4, total=13, industry='不動産（地域特化 戸建て販売）',
        role='顧客担当 兼 案件統括',
        channels='リスティング / ディスプレイ / SNS / 認知 / LINE / IG / WP制作',
        challenge='単一チャネル施策では地域内シェア獲得に限界。マーケティング全体の統合的な設計が必要だった。',
        approach='広告・SNS・LINE・インフルエンサー・WP制作までを横断的にディレクション。チャネル横断の統合マーケ施策を設計・実行。',
        kpis=[('145%', '来店数', GREEN), ('6ch', '統合チャネル数', PURPLE)],
        banner='地域特化 × 全6チャネル統合  ⇒  来店数 145%',
        diagram_fn=diag_case04,
    ),
    # BtoB / 人材グループ
    dict(
        num=5, total=13, industry='BtoB（知財系サービス）',
        role='顧客担当',
        channels='リスティング / SNS広告',
        challenge='業界特有の検索クエリ構造を把握しないと費用が無駄打ちになる難領域。一般的な運用ノウハウでは通用しなかった。',
        approach='業界特有の検索市場構造を深く分析し、Tier別のキャンペーン設計を構築。市場特性を最大限活かした運用設計に落とし込み、高CVR領域を最優先化。',
        kpis=[('社内', 'リード/契約数 ギネス', AMBER), ('独自', '市場構造分析', BLUE)],
        banner='広告経由リード数・契約数 ともに 社内ギネスを達成',
        diagram_fn=diag_case05,
    ),
    dict(
        num=6, total=13, industry='人材（特化型転職エージェント）',
        role='顧客担当 兼 案件統括',
        channels='リスティング / SNS / LPO',
        challenge='前体制下でCPAが高止まり。獲得効率の構造的な改善が必要で、広告だけでは限界があった。',
        approach='戦略設計から実行まで一気通貫で推進。自社アセットを活用したLP改善を主導し、CVR 125%を実現。プロジェクト統括も担当。',
        kpis=[('-50%', '登録CPA（3ヶ月）', GREEN), ('125%', 'CVR', GREEN)],
        banner='登録CPA -50%  |  CVR 125%  |  プロジェクト推進・統括も担当',
        diagram_fn=diag_case06,
    ),
    dict(
        num=7, total=13, industry='BtoB（SNS拡張ツール）',
        role='運用担当',
        channels='リスティング / SNS / サイト改善・LPO',
        challenge='既存の運用設計が現状の顧客像とずれており、商談数・効率ともに頭打ち。広告KPIだけでは課題が見えなかった。',
        approach='顧客ヒアリング・商談データ・Salesforceデータを横断的に分析し、運用設計を根本から再構築。サイト改善・LPOも同時に推進。',
        kpis=[('1.5倍', '商談数', GREEN), ('135%', 'CVR', GREEN)],
        banner='広告領域に閉じず、データ起点で運用設計を再構築する手法',
        diagram_fn=diag_case07,
    ),
    dict(
        num=8, total=13, industry='BtoB（領域特化 採用サービス）',
        role='顧客担当',
        channels='リスティング / DSP / SNS / 営業支援 / DM / AI活用',
        challenge='広告施策だけではリード数の拡大に限界。営業体制側にも商談化のボトルネックが存在していた。',
        approach='広告領域を超え、営業体制の課題を特定。営業代行施策を企画・実行し、商談獲得のフローを根本から強化。',
        kpis=[('+130%', 'リード数', GREEN), ('+50件/月', '営業代行 商談', GREEN)],
        banner='リード 130%  |  商談 120%  |  営業代行で新規 +50件/月',
        diagram_fn=diag_case08,
    ),
    # EC グループ
    dict(
        num=9, total=13, industry='EC（住宅設備・建設資材）',
        role='顧客担当',
        channels='リスティング / DSP / SNS / ダイナミック / Instagram運用',
        challenge='ECで売上拡張を目指すには刈り取り型施策だけでは限界。認知から購入までのフルファネル設計が必要だった。',
        approach='潜在〜顕在までのフルファネル施策を企画・提案から実行まで一気通貫で担当。媒体ごとの役割を再定義し、運用費の拡大投資を獲得。',
        kpis=[('260%', '担当領域経由 売上', GREEN), ('約4倍', '運用費 拡大', AMBER)],
        banner='担当領域経由の売上 260%  |  運用費 約4倍に拡大獲得',
        diagram_fn=diag_case09,
    ),
    # 店舗系 / 買取系 グループ
    dict(
        num=10, total=13, industry='店舗系（車整備）',
        role='顧客担当 兼 案件統括',
        channels='リスティング広告',
        challenge='全国展開の店舗。エリアによって需要構造が大きく異なり、画一的な配信ではCPAが安定しなかった。',
        approach='車保有人口・高速道路有無・店舗周囲の地形等の地域特性データを独自調査・分析。エリアごとの配信戦略を精緻に設計し直し。',
        kpis=[('-25%', 'CPA', GREEN), ('独自', 'エリア配信戦略', AMBER)],
        banner='独自データ調査 × 地域戦略設計  ⇒  CPA 25%改善',
        diagram_fn=diag_case10,
    ),
    dict(
        num=11, total=13, industry='店舗系（歯科矯正）',
        role='顧客担当',
        channels='リスティング / SNS広告',
        challenge='動画クリエイティブ制作で出演モデルが確保できず、クリエイティブ量産が制約に。改善サイクルが回らない状態だった。',
        approach='AI技術を活用して動画出演モデルを生成。クリエイティブ制作のボトルネックをテクノロジーで突破し、量産・高速ABテストを実現。',
        kpis=[('1/3', '予約CPA', GREEN), ('-50%', '来院CPA', GREEN)],
        banner='予約CPA  1/3  |  来院CPA  約-50%   ⇒  テクノロジーで制作の限界を突破',
        diagram_fn=diag_case11,
    ),
    dict(
        num=12, total=13, industry='店舗系（着物）',
        role='顧客担当 兼 案件統括',
        channels='リスティング / DSP / LINE運用',
        challenge='定量データだけのLP改善ではCVR改善が頭打ち。ユーザーが「なぜ離脱するのか」が見えていなかった。',
        approach='定量データに加え、定性的なユーザー行動データ（ヒートマップ・セッション）も活用したLP改善仮説検証を設計・実施。改善サイクルを構築。',
        kpis=[('約150%', 'CVR（対象ページ）', GREEN), ('構築', '改善サイクル', BLUE)],
        banner='対象ページ CVR  約150%改善  |  データドリブンな改善サイクルを構築',
        diagram_fn=diag_case12,
    ),
    dict(
        num=13, total=13, industry='買取系（商業車）',
        role='運用担当',
        channels='リスティング / ディスプレイ',
        challenge='問合せから買取への転換率が安定せず、単価と運用費のバランスが取れない状況。商材特性も独自性が高かった。',
        approach='市場・商材特性を深く分析した運用設計を構築。クリエイティブ・地域・デバイスを多角的に最適化し、転換率向上を実現。',
        kpis=[('+20%', '買取単価', GREEN), ('約3倍', '運用費 拡大', AMBER)],
        banner='市場・商材分析  ×  多軸最適化  ⇒  買取単価 +20%  |  運用費 約3倍',
        diagram_fn=diag_case13,
    ),
]

# === Group dividers ===
GROUPS = [
    dict(
        title='不動産', subtitle='REAL  ESTATE  /  4 CASES',
        accent=BLUE, accent_lt=BLUE_LT, icon=icon_realestate,
        cases=[(1, '不動産（売買領域）大規模運用', '問合せ +130%  |  月2億円規模'),
               (2, '不動産（法人向け売買）', '問合せ +140%  |  為替連動ロジック'),
               (3, '不動産（賃貸管理）', '問合せ +120%  |  新規立ち上げ'),
               (4, '不動産（地域特化 戸建て）', '来店 145%  |  全6チャネル統合')]
    ),
    dict(
        title='BtoB / 人材', subtitle='B2B  &  HR  /  4 CASES',
        accent=PURPLE, accent_lt=PURPLE_LT, icon=icon_btob,
        cases=[(5, 'BtoB（知財系サービス）', '社内ギネス達成  |  市場構造分析'),
               (6, '人材（特化型転職エージェント）', 'CPA -50%  |  CVR 125%'),
               (7, 'BtoB（SNS拡張ツール）', '商談 1.5倍  |  データ横断分析'),
               (8, 'BtoB（領域特化 採用サービス）', 'リード +130%  |  営業代行 +50/月')]
    ),
    dict(
        title='EC', subtitle='E-COMMERCE  /  1 CASE',
        accent=AMBER, accent_lt=AMBER_LT, icon=icon_ec,
        cases=[(9, 'EC（住宅設備・建設資材）', '売上 260%  |  運用費 約4倍 拡大')]
    ),
    dict(
        title='店舗系 / 買取', subtitle='RETAIL  /  4 CASES',
        accent=GREEN, accent_lt=GREEN_LT, icon=icon_shop,
        cases=[(10, '店舗系（車整備）', 'CPA -25%  |  独自エリア戦略'),
               (11, '店舗系（歯科矯正）', '予約CPA 1/3  |  AI動画モデル'),
               (12, '店舗系（着物）', 'CVR 約150%  |  定量×定性 改善'),
               (13, '買取系（商業車）', '単価 +20%  |  運用費 約3倍')]
    ),
]

# Map case num → group index
NUM_TO_GROUP = {}
for gi, g in enumerate(GROUPS):
    for c in g['cases']:
        NUM_TO_GROUP[c[0]] = gi


# === Main: open, build new slides, reorder, delete old ===
def main():
    prs = Presentation(SRC)

    # blank layout (no title placeholders) — typically layout index 6 in default theme
    blank_layout = None
    for layout in prs.slide_layouts:
        if 'Blank' in layout.name or 'blank' in layout.name or layout.name == '白紙':
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]

    # Create new slides at end. We'll then reorder.
    new_slide_ids = []  # list of (kind, group_idx_or_case_num, slide_obj_idx)

    sequence = []  # ordered list of slides to insert
    for gi, group in enumerate(GROUPS):
        # divider
        s = prs.slides.add_slide(blank_layout)
        # remove default placeholders if any
        _strip_placeholders(s)
        make_divider(s, group['title'], group['subtitle'],
                     group['cases'], group['accent'], group['accent_lt'], group['icon'])
        sequence.append(s)
        # cases in this group
        for c in group['cases']:
            num = c[0]
            # find case data
            case = next(cc for cc in CASES if cc['num'] == num)
            cs = prs.slides.add_slide(blank_layout)
            _strip_placeholders(cs)
            make_case_slide(cs, case)
            sequence.append(cs)

    # Reorder: keep slides 1, 2, 3 (idx 0,1,2), then sequence, then service menu (orig idx 6)
    # Original slides: 0=cover,1=profile,2=strengths,3=case1/3,4=case2/3,5=case3/3,6=menu
    # After adding new: still 0-6 original + new (sequence) appended.
    # Target order:
    # [0 cover, 1 profile, 2 strengths, ...sequence..., 6 menu] and DELETE 3,4,5.

    xml_slides = prs.slides._sldIdLst
    slide_elements = list(xml_slides)

    # snapshot original positions
    # indices: 0..6 original, 7..7+len(sequence)-1 new
    orig_keep = [slide_elements[0], slide_elements[1], slide_elements[2]]
    menu_slide = slide_elements[6]
    new_slides = slide_elements[7:]

    # remove all from xml_slides
    for el in slide_elements:
        xml_slides.remove(el)

    # add in new order: keep + new + menu
    for el in orig_keep + new_slides + [menu_slide]:
        xml_slides.append(el)

    # Now delete orphan slide parts (the 3 deleted slides 3,4,5)
    # Their slide objects are still in the package but not referenced.
    # We need to drop the relationships from the presentation part.
    pres_part = prs.part
    orphan_ids = [slide_elements[3].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'),
                  slide_elements[4].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'),
                  slide_elements[5].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')]
    # drop rels
    for rid in orphan_ids:
        try:
            pres_part.rels.pop(rid)
        except Exception:
            pass

    prs.save(DST)
    print(f'保存: {DST}')
    # Print final order
    prs2 = Presentation(DST)
    print(f'最終スライド数: {len(prs2.slides)}')
    for i, s in enumerate(prs2.slides):
        title = ''
        for shp in s.shapes:
            if shp.has_text_frame:
                t = shp.text_frame.text.strip()
                if t:
                    title = t.split('\n')[0][:50]
                    break
        print(f'  {i+1:02d}: {title}')


def _strip_placeholders(slide):
    """Remove default title/content placeholders from the slide."""
    sp_tree = slide.shapes._spTree
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    for sp in sp_tree.findall('.//p:sp', nsmap):
        nvSpPr = sp.find('p:nvSpPr', nsmap)
        if nvSpPr is not None:
            nvPr = nvSpPr.find('p:nvPr', nsmap)
            if nvPr is not None and nvPr.find('p:ph', nsmap) is not None:
                sp_tree.remove(sp)


if __name__ == '__main__':
    main()
